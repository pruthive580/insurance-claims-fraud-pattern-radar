"""Render the real HTML decks to pixel-perfect PPTX + PDF by screenshotting each
slide with Playwright's own Chromium, then packing the images full-bleed.
Run: .venv/bin/python render_decks.py
"""
import pathlib, tempfile
from playwright.sync_api import sync_playwright
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

WEB = pathlib.Path("fraud_radar/web").resolve()
DECKS = [
    ("presentation.html", "../Fraud-Pattern-Radar-Overwatch"),
    ("architecture.html", "../Fraud-Pattern-Radar-Architecture"),
]

def render(html_name, out_base):
    src = (WEB / html_name).as_uri()
    tmp = pathlib.Path(tempfile.mkdtemp())
    shots = []
    with sync_playwright() as p:
        b = p.chromium.launch()
        pg = b.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=2)
        pg.goto(src)
        pg.wait_for_timeout(500)
        # Kill entry fade/animations so no slide is ever captured mid-transition.
        pg.add_style_tag(content=".slide{transition:none !important;animation:none !important}"
                                 " .draw,.pulse{animation:none !important}")
        n = pg.evaluate("document.querySelectorAll('.slide').length")
        for i in range(n):
            pg.evaluate(f"go({i})")
            pg.wait_for_timeout(350)          # SVG paint settle (no transition to wait on now)
            f = tmp / f"s{i:02}.png"
            pg.screenshot(path=str(f))         # current viewport = this slide
            shots.append(f)
        b.close()

    imgs = [Image.open(f).convert("RGB") for f in shots]
    imgs[0].save(out_base + ".pdf", save_all=True, append_images=imgs[1:])

    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for f in shots:
        sl = prs.slides.add_slide(blank)
        sl.shapes.add_picture(str(f), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(out_base + ".pptx")
    return n, imgs[0].size

for html, base in DECKS:
    n, size = render(html, base)
    print(f"{html}: {n} slides @ {size[0]}x{size[1]}px  ->  {base}.pptx + .pdf")
