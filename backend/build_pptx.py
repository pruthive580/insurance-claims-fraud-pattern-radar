"""Generate an editable PowerPoint of the Fraud Pattern Radar pitch deck.
Native shapes/text (not screenshots) so the team can edit it in PowerPoint.
Run: .venv/bin/python build_pptx.py  ->  ../Fraud-Pattern-Radar-Overwatch.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

C = dict(bg="0b0f16", card="10151d", cardhi="0f1a27", line="2a3340", ink="e6edf3",
         muted="8b949e", accent="58a6ff", accent2="a371f7", green="3fb950",
         high="f85149", med="d29922", low="2ea043", gtext="5bd67c", ptext="c8a6ff",
         redcard="2a1417", greencard="12241a")
def rgb(h): return RGBColor.from_string(h)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb(C["bg"])
    return s

def _runs(p, runs, align):
    p.alignment = align
    for txt, sz, col, bold in runs:
        r = p.add_run(); r.text = txt
        f = r.font; f.size = Pt(sz); f.bold = bold; f.color.rgb = rgb(col); f.name = "Aptos"

def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    if isinstance(runs, tuple): runs = [runs]
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    _runs(tf.paragraphs[0], runs, align)
    return box

def paras(s, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, runs in enumerate(lines):
        if isinstance(runs, tuple): runs = [runs]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space); _runs(p, runs, align)
    return box

def rrect(s, l, t, w, h, fill=C["card"], line=C["line"], lw=1.0, lines=None, anchor=MSO_ANCHOR.MIDDLE):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
    if line: shp.line.color.rgb = rgb(line); shp.line.width = Pt(lw)
    else: shp.line.fill.background()
    shp.shadow.inherit = False
    if lines:
        tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        for i, runs in enumerate(lines):
            if isinstance(runs, tuple): runs = [runs]
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _runs(p, runs, PP_ALIGN.CENTER)
    return shp

def oval(s, l, t, w, h, fill, line, lw=2.0, lines=None):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line); shp.line.width = Pt(lw); shp.shadow.inherit = False
    if lines:
        tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, runs in enumerate(lines):
            if isinstance(runs, tuple): runs = [runs]
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _runs(p, runs, PP_ALIGN.CENTER)
    return shp

def arrow(s, x1, y1, x2, y2, color=C["accent"], w=2.0, dash=False):
    cxn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cxn.line.color.rgb = rgb(color); cxn.line.width = Pt(w)
    ln = cxn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if dash:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    return cxn

def header(s, eyebrow, title, tcolor=C["ink"]):
    text(s, 0.7, 0.45, 12, 0.4, (eyebrow.upper(), 13, C["accent"], True))
    text(s, 0.7, 0.86, 12, 1.0, (title, 28, tcolor, True))

def tag(s, l, t, txt):
    return rrect(s, l, t, 0.1 + 0.11 * len(txt), 0.34, fill=C["bg"], line=C["line"],
                 lines=[[(txt, 11, C["muted"], False)]])

# ---------------------------------------------------------------- slides
# 1 TITLE
s = slide()
text(s, 0.9, 1.9, 11.5, 0.5, ("HACKATHON 2026 · USE CASE #12 · INSURANCE / SIU", 14, C["accent"], True))
paras(s, 0.9, 2.4, 11.5, 2.2, [
    [("Fraud Pattern Radar", 46, C["ink"], True)],
    [("Catch the whole ring — and prove it.", 34, C["accent"], True)],
], space=4)
text(s, 0.9, 4.5, 10.8, 1.2, ("Explainable, real-time claim-fraud scoring with organised-ring detection. "
     "Every flag traces to a named rule — the maths decides; AI just helps at the edges.", 18, "cdd9e5", False))
for i, tg in enumerate(["Team Overwatch", "Deterministic-first", "Network-aware", "$0 local AI"]):
    tag(s, 0.9 + i * 2.5, 6.0, tg)

# 2 PROBLEM
s = slide(); header(s, "The problem", "Fraud hides between claims, not inside one")
cards = [("Rings are invisible", "Three ‘unrelated’ claimants can share a phone, a bank account and a repair shop — a single-claim review never sees it."),
         ("Adjusters are buried", "Volume forces fast calls. Padded and staged claims slip through; the real fraud waits in the queue."),
         ("Black-box scores don’t fly", "SIU and regulators need a defensible reason for every referral — not an opaque probability.")]
for i, (h, b) in enumerate(cards):
    x = 0.7 + i * 4.13
    r = rrect(s, x, 2.1, 3.9, 2.9, fill=C["card"], anchor=MSO_ANCHOR.TOP)
    paras(s, x + 0.25, 2.35, 3.4, 2.5, [[(h, 18, C["ink"], True)], [(b, 13.5, C["muted"], False)]], space=8)
text(s, 0.7, 5.6, 12, 0.6, [("Insurers need scoring that is ", 15, C["muted"], False),
     ("fast at intake, explainable, and network-aware.", 15, C["ink"], True)])

# 3 ARCHITECTURE
s = slide(); header(s, "Architecture", "Small, layered — deterministic core, dataset-driven")
rrect(s, 0.7, 2.3, 2.5, 2.4, fill=C["cardhi"], line=C["accent"], lw=1.6,
      lines=[[("Browser SPA", 15, C["ink"], True)], [("queue · ring graph", 12, C["muted"], False)],
             [("why-flagged · FNOL", 12, C["muted"], False)]])
rrect(s, 3.9, 2.55, 2.1, 1.9, lines=[[("FastAPI", 15, C["ink"], True)], [("REST · Pydantic v2", 12, C["muted"], False)]])
core = rrect(s, 6.7, 1.9, 3.4, 3.0, fill=C["greencard"], line=C["green"], lw=2.4,
             lines=[[("DETERMINISTIC CORE", 15, C["gtext"], True)], [("scoring.py orchestrator", 12, C["muted"], False)],
                    [("rules.py · R1–R16", 13, C["ink"], True)], [("graph.py · rings", 13, C["ink"], True)],
                    [("llm.py · 🧪 Rule Lab", 13, C["ptext"], True)]])
# dataset cylinder
cyl = s.shapes.add_shape(MSO_SHAPE.CAN, Inches(10.7), Inches(2.4), Inches(1.9), Inches(2.0))
cyl.fill.solid(); cyl.fill.fore_color.rgb = rgb(C["card"]); cyl.line.color.rgb = rgb(C["line"]); cyl.shadow.inherit = False
text(s, 10.7, 4.5, 1.9, 0.6, [("claims.json", 12, C["ink"], True)], align=PP_ALIGN.CENTER)
arrow(s, 3.2, 3.5, 3.85, 3.5); arrow(s, 6.0, 3.5, 6.65, 3.5); arrow(s, 10.1, 3.4, 10.65, 3.4)
arrow(s, 6.7, 4.6, 3.2, 4.6, color=C["green"], dash=True)
text(s, 3.4, 4.65, 3.2, 0.4, [("score · band · reasons · ring", 11, C["gtext"], False)], align=PP_ALIGN.CENTER)
text(s, 0.7, 5.7, 12, 0.8, [("Scoring is ", 14, C["muted"], False), ("100% reproducible", 14, C["ink"], True),
     (" — the AI works only at the edges (reads claims, writes notes/referrals, proposes rules) on ", 14, C["muted"], False),
     ("local Ollama", 14, C["ink"], True), (". The core never depends on it.", 14, C["muted"], False)])

# 4 PIPELINE
s = slide(); header(s, "How it works", "One claim → a routed, explained decision")
steps = [("Intake / FNOL", "API or console"), ("Context", "percentiles · thresholds"),
         ("Cross-claim", "rings · dup photos"), ("Rules R1–R16", "points + reason"),
         ("Band", "Low / Med / High"), ("Route + Note", "action · narrative")]
bw = 1.86
for i, (h, b) in enumerate(steps):
    x = 0.55 + i * (bw + 0.15)
    fill = C["cardhi"] if i == 3 else C["card"]
    ln = C["accent"] if i == 3 else C["line"]
    rrect(s, x, 2.6, bw, 1.5, fill=fill, line=ln, lines=[[(h, 13, C["ink"], True)], [(b, 11, C["muted"], False)]])
    if i < 5: arrow(s, x + bw, 3.35, x + bw + 0.15, 3.35)
text(s, 0.7, 4.6, 12, 0.5, [("Worked example (walk-in): ", 14, C["muted"], True),
     ("R1 +15 · R12 +15 · R13 +10 · R8 +12 · R7 +10 · R6 +5 · R11 +7  =  ", 14, C["ink"], False),
     ("74 → HIGH → SIU", 15, C["high"], True)])
text(s, 0.7, 5.5, 12, 0.6, [("The score is always the ", 14, C["muted"], False),
     ("exact sum of fired-rule points", 14, C["ink"], True), (" — additive, reproducible, auditable.", 14, C["muted"], False)])

# 5 THE 16 RULES
s = slide(); header(s, "The engine", "16 rules across 6 fraud categories")
RULES = [("R1", "Just-in-time cover", 15), ("R2", "Reinstated pre-loss", 10), ("R3", "Address mismatch", 8),
         ("R4", "Late report", 8), ("R5", "Serial claimant", 12), ("R6", "Weak-witness timing", 5),
         ("R7", "Above 90th pct amount", 10), ("R8", "Threshold gaming", 12), ("R9", "Fabricated invoices", 8),
         ("R10", "Fast cash, no inspect", 10), ("R11", "No police report", 7), ("R12", "Shared-identifier ring", 15),
         ("R13", "Flagged repairer", 10), ("R14", "Doc inconsistency", 10), ("R15", "Duplicate photo", 15),
         ("R16", "EXIF mismatch", 8)]
cw, ch = 3.0, 0.86
for i, (rid, nm, pts) in enumerate(RULES):
    col, row = i % 4, i // 4
    x = 0.7 + col * (cw + 0.1); y = 2.2 + row * (ch + 0.14)
    rrect(s, x, y, cw, ch, lines=[[(f"{rid}  ", 13, C["accent"], True), (nm, 12.5, C["ink"], False),
                                    (f"   +{pts}", 12, C["high"], True)]])
text(s, 0.7, 6.35, 12, 0.5, [("Every rule tested positive, negative, and at its exact boundary — plus a human-in-the-loop Rule Lab to add more.", 13.5, C["muted"], False)])

# 6 BANDS
s = slide(); header(s, "Score bands → action", "The score drives a concrete routing decision")
rows = [("LOW", "0 – 29", "Straight-through / auto-adjudicate", C["low"]),
        ("MEDIUM", "30 – 59", "Adjuster review with reason codes", C["med"]),
        ("HIGH", "≥ 60", "Route to SIU queue — top of the list", C["high"])]
for i, (band, rng, act, col) in enumerate(rows):
    y = 2.5 + i * 1.1
    rrect(s, 0.7, y, 1.7, 0.8, fill=col, line=col, lines=[[(band, 15, C["bg"], True)]])
    text(s, 2.6, y, 1.8, 0.8, [(rng, 18, C["ink"], True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.6, y, 8.0, 0.8, [(act, 17, C["ink"], False)], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.7, 6.1, 12, 0.6, [("Straight from the #12/#48 rule card — this one number powers the ranked queue and the auto-routing.", 14, C["muted"], False)])

# 7 RING (entity graph)
s = slide(); header(s, "The centrepiece", "Fraud-ring detection via network analysis")
# claimant circles
oval(s, 1.6, 2.2, 1.15, 1.15, C["redcard"], C["high"], lines=[[("Alice", 12, "ff7b72", True)], [("C-101 · 82", 10, C["muted"], False)]])
oval(s, 1.6, 4.6, 1.15, 1.15, C["redcard"], C["high"], lines=[[("Bruno", 12, "ff7b72", True)], [("C-102 · 69", 10, C["muted"], False)]])
oval(s, 10.4, 3.4, 1.15, 1.15, C["redcard"], C["high"], lines=[[("Carla", 12, "ff7b72", True)], [("C-103 · 75", 10, C["muted"], False)]])
# shared hubs
hubs = [("📞 555-0100", "shared phone", 2.35), ("🔧 QuickFix", "same repairer", 3.55), ("⚖️ Shady & Co", "same lawyer", 4.75)]
for txt, sub, y in hubs:
    rrect(s, 5.7, y, 1.9, 0.95, fill=C["cardhi"], line=C["accent"], lines=[[(txt, 12.5, C["ink"], True)], [(sub, 10, C["muted"], False)]])
# edges (approx center-to-center)
CA, CB, CC = (2.75, 2.78), (2.75, 5.18), (10.4, 3.98)
for cx, cy in (CA, CB, CC):
    for hy in (2.83, 4.03, 5.23):
        arrow(s, cx, cy, 5.7 if cx < 5 else 7.6, hy, color="3b4553", w=1.5)
text(s, 0.7, 6.25, 12, 0.7, [("Three different claimants, three policies — tied together by a shared phone, shop and lawyer. ", 14, C["muted"], False),
     ("That shared box IS the ring.", 14, "ff7b72", True)])

# 8 FNOL LIVE
s = slide(); header(s, "Live at intake", "A brand-new claim joins the ring in real time")
rrect(s, 0.7, 2.3, 5.6, 3.0, fill=C["card"], anchor=MSO_ANCHOR.TOP,
      lines=[[("The walk-in", 18, C["ink"], True)],
             [("Never-seen claim, scored on submission.", 13.5, C["muted"], False)],
             [("Reuses the ring’s phone + repairer.", 13.5, C["muted"], False)]])
rrect(s, 0.95, 3.7, 5.1, 1.3, fill=C["greencard"], line=C["high"], lw=1.5,
      lines=[[("74", 40, C["high"], True), ("   HIGH → SIU · RING-1", 16, C["ink"], True)]])
paras(s, 6.6, 2.3, 6.0, 3.0, [
    [("Why (live)", 18, C["ink"], True)],
    [("R12 +15 — shares phone & repairer with C-101/102/103", 14, "cdd9e5", False)],
    [("R13 +10 — repairer has prior SIU flags", 14, "cdd9e5", False)],
    [("R1 +15 — filed 3 days after inception", 14, "cdd9e5", False)],
    [("R8 +12 — $24.5k just under the $25k threshold", 14, "cdd9e5", False)],
], space=10)
text(s, 0.7, 5.7, 12, 0.6, [("In the dashboard a node animates into the graph and links to RING-1 — detection is ", 14, C["muted"], False),
     ("dynamic, not precomputed.", 14, C["ink"], True)])

# 9 AI AT THE EDGES (+ $0)
s = slide(); header(s, "AI at the edges", "AI wraps the engine — it never decides the verdict")
rrect(s, 0.7, 2.2, 2.4, 1.1, fill=C["card"], lines=[[("📄 Messy claim", 13, C["ink"], True)], [("free text · PDF", 11, C["muted"], False)]])
arrow(s, 3.1, 2.75, 3.9, 2.75, color=C["accent2"])
rrect(s, 3.95, 2.0, 4.0, 1.5, fill=C["greencard"], line=C["green"], lw=2.2,
      lines=[[("🛡️ DETERMINISTIC CORE", 14, C["gtext"], True)], [("score · band · ring — AI-free", 11.5, C["muted"], False)]])
arrow(s, 7.95, 2.75, 8.75, 2.75, color=C["accent2"])
rrect(s, 8.8, 2.05, 3.8, 1.4, fill=C["card"],
      lines=[[("🧠 note · 🕸️ ring summary", 12.5, C["ink"], True)], [("📝 SIU referral", 12.5, C["ink"], True)]])
rrect(s, 3.95, 3.75, 4.0, 0.7, fill="1a1526", line=C["accent2"], lw=1.4,
      lines=[[("🧪 Rule Lab: AI proposes → human adopts", 11.5, C["ptext"], True)]])
# $0 highlight banner
rrect(s, 0.7, 4.9, 11.9, 1.2, fill=C["greencard"], line=C["green"], lw=1.8,
      lines=[[("💵 $0 to run", 22, C["gtext"], True), ("   — local Ollama · no per-token fees · data never leaves your walls", 16, C["ink"], True)]])
text(s, 0.7, 6.4, 12, 0.5, [("Messy → structured in, structured → readable out — the number in the middle stays pure maths.", 13.5, C["muted"], False)])

# 10 COVERAGE
s = slide(); header(s, "Scope coverage", "All 6 use cases · all 4 MVP items")
cov = [("1 · Scoring at FNOL / intake", "Live POST /api/score + in-UI console"),
       ("2 · Anomaly: amount / timing / frequency", "R7,R8 · R1,R2,R4,R6 · R5"),
       ("3 · Entity linking", "Shared-identifier graph + R5 / R13"),
       ("4 · SIU prioritization", "Ranked queue + band routing (+ risk×$)"),
       ("5 · Explanation layer", "Per-rule reasons + narrative + additive score"),
       ("6 · Fraud-ring detection", "Connected-component network analysis")]
for i, (req, by) in enumerate(cov):
    y = 2.2 + i * 0.72
    text(s, 0.7, y, 5.6, 0.6, [(req, 14.5, C["ink"], True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.3, y, 5.7, 0.6, [(by, 13, C["muted"], False)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 12.1, y, 0.6, 0.6, [("✓", 16, C["gtext"], True)], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.7, 6.7, 12, 0.4, [("Every row backed by an automated conformance test — 103 checks, all green.", 13.5, C["muted"], False)])

# 11 QUALITY
s = slide(); header(s, "Quality", "Built to be trusted, not just demoed")
stats = [("103", "automated checks — all green", C["gtext"]), ("16 × 3", "rules × pos / neg / boundary", C["ink"]),
         ("404 · 422", "API robustness covered", C["ink"]), ("0 FP", "clean controls stay clean", C["gtext"])]
for i, (b, l, col) in enumerate(stats):
    x = 0.7 + i * 3.05
    rrect(s, x, 2.5, 2.85, 1.7, fill=C["card"], lines=[[(b, 30, col, True)], [(l, 12, C["muted"], False)]])
paras(s, 0.7, 4.7, 12, 1.6, [
    [("→ Every rule fires when it should — and stays silent when it shouldn’t (boundary-exact).", 15, "dfe7ef", False)],
    [("→ Ring logic verified: distinct-claimant requirement, edge de-dup, same-claimant ≠ ring.", 15, "dfe7ef", False)],
    [("→ Score additivity asserted: Σ points == score on every claim.", 15, "dfe7ef", False)],
], space=8)

# 12 TECH (+ $0 highlight)
s = slide(); header(s, "Tech", "Boring where it counts, sharp where it matters")
layers = [("🖥️ Presentation", "zero-build SPA · vanilla JS · light/dark", C["card"], C["line"]),
          ("🔌 API", "FastAPI · Pydantic v2 · REST", C["card"], C["line"]),
          ("🧠 Domain logic", "rule engine R1–R16 + graph · 100% reproducible", C["greencard"], C["green"]),
          ("🗂️ Data", "claims.json / config.json — swappable", C["card"], C["line"])]
for i, (h, b, fill, ln) in enumerate(layers):
    y = 2.2 + i * 0.82
    rrect(s, 0.7, y, 8.4, 0.72, fill=fill, line=ln, lw=(2.2 if i == 2 else 1.0),
          lines=[[(h + "   —   ", 14, C["ink"], True), (b, 12.5, C["muted"], False)]])
rrect(s, 9.4, 2.6, 3.2, 1.5, fill="1a1526", line=C["accent2"], lw=1.6,
      lines=[[("🤖 Local AI", 14, C["ptext"], True)], [("Ollama qwen3", 12, C["muted"], False)], [("intake · notes · Rule Lab", 11.5, C["muted"], False)]])
rrect(s, 9.4, 4.25, 3.2, 1.1, fill=C["greencard"], line=C["green"], lw=1.8,
      lines=[[("💵 $0 AI", 22, C["gtext"], True)], [("no per-token fees", 12, C["muted"], False)]])
text(s, 0.7, 5.7, 8.4, 0.8, [("~1,300 lines · 4 runtime deps · runs fully offline. A claim falls down the stack; local AI hangs off the side.", 13.5, C["muted"], False)])

# 13 ROADMAP + CLOSE
s = slide(); header(s, "What’s next · thank you", "Team Overwatch")
text(s, 0.7, 1.9, 12, 1.0, [("Catch the whole ring — and prove it.", 34, C["accent"], True)])
road = [("Fuzzy entity linking ✓", "normalises phones/addresses (shipped)"),
        ("Real media forensics", "EXIF / pHash on uploaded photos"),
        ("Official dataset", "drop in #12/#48 data — no code change")]
for i, (h, b) in enumerate(road):
    x = 0.7 + i * 4.13
    rrect(s, x, 3.1, 3.9, 1.7, fill=C["card"], anchor=MSO_ANCHOR.TOP,
          lines=[[(h, 15, C["ink"], True)], [(b, 12.5, C["muted"], False)]])
text(s, 0.7, 5.1, 12, 0.8, [("Explainable · network-aware · dataset-driven · offline-safe · $0 AI — mapped to every use case and MVP item.", 16, "cdd9e5", False)])
text(s, 0.7, 6.3, 12, 0.4, [("dashboard → localhost:8000   ·   103/103 tests green", 13, C["muted"], False)])

out = "../Fraud-Pattern-Radar-Overwatch.pptx"
prs.save(out)
print("saved", out, "·", len(prs.slides._sldIdLst), "slides")
